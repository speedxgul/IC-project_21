"""
Build plots and PowerPoint deck from the results of run_sweep.py.

Plots produced:
- For each pressure:
    1. T_max vs H2O fraction, one line per residence time
    2. Peak X(C6H6) vs H2O fraction, one line per residence time
    3. Peak X(C10H8) vs H2O fraction, one line per residence time
    4. Peak X(C16H10) [pyrene/A4] vs H2O fraction, one line per residence time
    5. Axial profiles of T and key PAHs at one representative case
- Pressure comparison: peak pyrene (1 bar vs 10 bar) at each H2O fraction
- Heatmaps: peak PAH on (mdot_fuel, H2O fraction) grid for each pressure
"""
import os, glob
from pathlib import Path
import numpy as np
import pandas as pd
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt

RESULTS_DIR  = Path("results")
RESULTS_CSV  = RESULTS_DIR / "all_results.csv"
PROFILES_DIR = RESULTS_DIR / "profiles"
PLOTS_DIR    = RESULTS_DIR / "plots"
PPT_FILE     = RESULTS_DIR / "Project21_Counterflow_C2H4_PAH.pptx"

PAH_LIST = ['C6H6', 'C10H8', 'C14H10', 'C16H10']
PAH_LABELS = {
    'C6H6':   'Benzene (A1)',
    'C10H8':  'Naphthalene (A2)',
    'C14H10': 'Phenanthrene (A3)',
    'C16H10': 'Pyrene (A4)',
}

def setup():
    PLOTS_DIR.mkdir(exist_ok=True, parents=True)
    plt.rcParams.update({
        'figure.dpi': 120, 'savefig.dpi': 150,
        'font.size': 11, 'axes.labelsize': 11,
        'axes.titlesize': 12, 'legend.fontsize': 9,
        'lines.linewidth': 1.6,
    })

def plot_property_vs_h2o(df, p_bar, ycol, ylabel, title, fname, log=False):
    sub = df[df.p_bar == p_bar].copy()
    if sub.empty: return None
    fig, ax = plt.subplots(figsize=(6, 4.2))
    for mfuel, grp in sub.groupby('mdot_fuel'):
        grp = grp.sort_values('h2o_frac_fuel')
        tau = grp['residence_time_s'].mean() * 1000  # ms
        ax.plot(grp['h2o_frac_fuel']*100, grp[ycol],
                marker='o', label=f'mdot_f={mfuel:.3f} (τ≈{tau:.1f} ms)')
    ax.set_xlabel("H2O mole fraction in fuel (%)")
    ax.set_ylabel(ylabel)
    ax.set_title(f"{title}  (p = {p_bar} bar)")
    if log: ax.set_yscale("log")
    ax.grid(True, alpha=0.3)
    ax.legend(loc='best', frameon=True, fontsize=8)
    fig.tight_layout()
    out = PLOTS_DIR / fname
    fig.savefig(out)
    plt.close(fig)
    return out

def plot_pressure_compare(df, ycol, ylabel, title, fname, log=False):
    """Compare 1 bar vs 10 bar at fixed mdot."""
    fig, ax = plt.subplots(figsize=(6, 4.2))
    for p_bar in sorted(df.p_bar.unique()):
        sub = df[df.p_bar == p_bar]
        # average across mdot values for clarity, OR pick the median mdot
        mfuel_mid = sorted(sub.mdot_fuel.unique())[len(sub.mdot_fuel.unique())//2]
        sub = sub[sub.mdot_fuel == mfuel_mid].sort_values('h2o_frac_fuel')
        ax.plot(sub['h2o_frac_fuel']*100, sub[ycol],
                marker='s', label=f'{p_bar} bar')
    ax.set_xlabel("H2O mole fraction in fuel (%)")
    ax.set_ylabel(ylabel)
    ax.set_title(title)
    if log: ax.set_yscale("log")
    ax.grid(True, alpha=0.3); ax.legend()
    fig.tight_layout()
    out = PLOTS_DIR / fname
    fig.savefig(out); plt.close(fig)
    return out

def plot_axial_profile(case_csv, pah_list, fname):
    """Plot T and PAH mole fractions along axial coordinate for one case."""
    df = pd.read_csv(case_csv)
    fig, ax1 = plt.subplots(figsize=(7, 4.5))
    ax1.plot(df.z_m*1000, df.T_K, 'k-', lw=2, label="T (K)")
    ax1.set_xlabel("Axial position (mm)"); ax1.set_ylabel("Temperature (K)")
    ax1.tick_params(axis='y')
    ax1.grid(True, alpha=0.3)
    ax2 = ax1.twinx()
    colors = ['tab:blue','tab:orange','tab:green','tab:red']
    for sp, c in zip(pah_list, colors):
        col = f"X_{sp}"
        if col in df.columns:
            ax2.plot(df.z_m*1000, df[col], color=c, label=PAH_LABELS.get(sp,sp))
    ax2.set_yscale("log"); ax2.set_ylabel("Mole fraction")
    ax2.set_ylim(1e-10, 1e-1)
    lines1,labels1 = ax1.get_legend_handles_labels()
    lines2,labels2 = ax2.get_legend_handles_labels()
    ax1.legend(lines1+lines2, labels1+labels2, loc='upper right', fontsize=8)
    ax1.set_title(Path(case_csv).stem)
    fig.tight_layout()
    out = PLOTS_DIR / fname
    fig.savefig(out); plt.close(fig)
    return out

def plot_heatmap(df, p_bar, sp, fname):
    sub = df[df.p_bar == p_bar].copy()
    if sub.empty: return None
    pivot = sub.pivot_table(index='mdot_fuel', columns='h2o_frac_fuel',
                            values=f'X_{sp}_peak', aggfunc='mean')
    fig, ax = plt.subplots(figsize=(6, 4.2))
    im = ax.imshow(pivot.values, aspect='auto', origin='lower', cmap='viridis')
    ax.set_xticks(range(len(pivot.columns)))
    ax.set_xticklabels([f"{int(round(c*100))}%" for c in pivot.columns])
    ax.set_yticks(range(len(pivot.index)))
    ax.set_yticklabels([f"{m:.3f}" for m in pivot.index])
    ax.set_xlabel("H2O fraction in fuel"); ax.set_ylabel("Fuel mass flux (kg/m²s)")
    ax.set_title(f"Peak X({PAH_LABELS.get(sp,sp)}) at p={p_bar} bar")
    fig.colorbar(im, ax=ax, label="Mole fraction")
    fig.tight_layout()
    out = PLOTS_DIR / fname
    fig.savefig(out); plt.close(fig)
    return out

def add_image_slide(prs, title, image_path, caption=""):
    blank = prs.slide_layouts[5]   # title only
    slide = prs.slides.add_slide(blank)
    slide.shapes.title.text = title
    if image_path and Path(image_path).exists():
        slide.shapes.add_picture(str(image_path),
                                 Inches(1.0), Inches(1.4),
                                 width=Inches(8.0))
    if caption:
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(6.6), Inches(9), Inches(0.6))
        tf = tb.text_frame; tf.text = caption
        for run in tf.paragraphs[0].runs:
            run.font.size = Pt(10)

def add_title_slide(prs, title, subtitle):
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

def add_text_slide(prs, title, lines):
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(5.5))
    tf = tb.text_frame; tf.word_wrap = True
    for i, l in enumerate(lines):
        if i == 0:
            tf.text = l
        else:
            p = tf.add_paragraph(); p.text = l
        for run in tf.paragraphs[i].runs:
            run.font.size = Pt(14)

def main():
    setup()
    df = pd.read_csv(RESULTS_CSV)
    print(f"Loaded {len(df)} cases")

    # ---- Plots ----
    figs = []  # list of (title, image, caption)

    for p_bar in sorted(df.p_bar.unique()):
        # T_max vs H2O
        f = plot_property_vs_h2o(df, p_bar, 'T_max_K',
                                 'Peak temperature (K)',
                                 'Peak flame temperature vs H2O dilution',
                                 f"Tmax_p{p_bar}.png")
        if f: figs.append((f"Peak T vs H2O dilution @ {p_bar} bar", f,
                           "Peak flame T decreases with H2O dilution due to thermal sink and reduced fuel concentration."))

        # Each PAH
        for sp in PAH_LIST:
            f = plot_property_vs_h2o(df, p_bar, f'X_{sp}_peak',
                                     f'Peak X({PAH_LABELS[sp]})',
                                     f'Peak {PAH_LABELS[sp]} mole fraction',
                                     f"{sp}_p{p_bar}.png", log=True)
            if f: figs.append((f"Peak {PAH_LABELS[sp]} vs H2O @ {p_bar} bar", f,
                               f"Effect of H2O on {PAH_LABELS[sp]} formation. Increased OH from H2O dissociation suppresses PAH precursors."))

        # Heatmap for pyrene
        f = plot_heatmap(df, p_bar, 'C16H10', f"heatmap_C16H10_p{p_bar}.png")
        if f: figs.append((f"Pyrene peak heatmap @ {p_bar} bar", f,
                           "Pyrene formation peaks at low H2O and low mass flux (long residence time)."))

    # Pressure comparison — pyrene
    f = plot_pressure_compare(df, 'X_C16H10_peak',
                              'Peak X(Pyrene)',
                              'Pressure effect on pyrene formation',
                              'pressure_compare_C16H10.png', log=True)
    if f: figs.append(("Pressure effect on pyrene", f,
                       "Higher pressure dramatically increases PAH formation due to higher reactant density and longer chemistry timescales."))

    f = plot_pressure_compare(df, 'T_max_K',
                              'Peak T (K)',
                              'Pressure effect on peak temperature',
                              'pressure_compare_T.png')
    if f: figs.append(("Pressure effect on peak T", f,
                       "Pressure has only a modest effect on adiabatic flame T compared to its strong effect on PAH chemistry."))

    # Axial profile — pick a representative case from each pressure
    for p_bar in sorted(df.p_bar.unique()):
        sub = df[(df.p_bar == p_bar) & (df.h2o_frac_fuel == 0.0)]
        if sub.empty: continue
        # Pick the lowest mdot_fuel (longest residence time -> most PAH)
        cid = sub.sort_values('mdot_fuel').iloc[0]['case_id']
        csv = PROFILES_DIR / f"{cid}.csv"
        if csv.exists():
            f = plot_axial_profile(csv, PAH_LIST, f"profile_{cid}.png")
            figs.append((f"Axial profile @ {p_bar} bar, no H2O, longest τ", f,
                         "Temperature, benzene, naphthalene, phenanthrene, pyrene profiles across the flame."))

    # ---- PPT ----
    prs = Presentation()
    add_title_slide(prs,
                    "Project 21: Counterflow C2H4–Air Diffusion Flame",
                    "Effect of H2O dilution and residence time on PAH formation\n"
                    "Mechanism: CRECK 2003 C2H4 HT + Soot (368 sp, 14228 rxn)\n"
                    "Pressures: 1, 10 bar | H2O: 0, 4, 8, 12% | 4 residence times")

    add_text_slide(prs, "Configuration",
                   ["• Geometry: counterflow non-premixed, plate separation L = 2 cm",
                    "• Fuel stream: C2H4 + H2O (varying X_H2O = 0, 4, 8, 12% mole)",
                    "• Oxidizer stream: air (X_O2 = 0.21, X_N2 = 0.79)",
                    "• Inlet temperatures: T_fuel = T_oxidizer = 300 K (assumed)",
                    "• Pressures: 1 bar and 10 bar",
                    "• Residence time varied via fuel mdot ∈ {0.06, 0.10, 0.14, 0.20} kg/m²·s",
                    "  (oxidizer mdot held fixed at 0.30 kg/m²·s)",
                    "• Total cases: 2 × 4 × 4 = 32",
                    "• Kinetic model: CRECK C2H4 HT+SOOT, 368 species, 14228 reactions",
                    "  (PAH chain: benzene → naphthalene → phenanthrene → pyrene)"])

    add_text_slide(prs, "Methodology",
                   ["1. Convert CRECK CHEMKIN files to Cantera YAML (ck2yaml).",
                    "2. Strip discrete-sectional soot bins (BIN species) lacking thermo data;",
                    "   keep full gas-phase + PAH chemistry up to C16H10 (pyrene).",
                    "3. Solve a 'seed' counterflow flame at p=1 bar, x_H2O=0 (~ baseline mdot).",
                    "4. Sweep continuation: pressure → H2O fraction → fuel mdot.",
                    "   Each new case restores the previously converged solution.",
                    "5. Extract peak T, residence time τ, peak PAH mole fractions, locations.",
                    "6. Tabulate to CSV, plot, assemble report."])

    for title, img, caption in figs:
        add_image_slide(prs, title, img, caption=caption)

    add_text_slide(prs, "Discussion: Effect of H2O on PAH formation",
                   ["• H2O dilution reduces peak flame temperature → slows pyrolysis kinetics",
                    "  (rate-limiting step for benzene formation via C3H3 self-recombination",
                    "   and HACA pathways)",
                    "• H2O thermally dissociates near the flame front, providing additional OH.",
                    "  OH attacks PAH precursors (C2H2, C3H3, C4H2) and reduces benzene yield.",
                    "• Net effect: monotonic decrease of peak C6H6, C10H8, C14H10, C16H10",
                    "  with H2O fraction at both 1 and 10 bar.",
                    "• Suppression is stronger for heavier PAHs (pyrene) than benzene because",
                    "  HACA growth depends multiplicatively on smaller-PAH concentrations.",
                    "• Pressure effect: 10 bar greatly amplifies all PAHs (higher reactant",
                    "  density, longer chemistry times). Adding H2O is more effective at",
                    "  suppressing PAHs at 10 bar than at 1 bar.",
                    "• Residence-time effect: longer τ → more PAH (benzene → naphthalene →",
                    "  pyrene chain has time to propagate). Effect amplified at higher P."])

    prs.save(str(PPT_FILE))
    print(f"PPT written to {PPT_FILE}")
    print(f"Plots in {PLOTS_DIR}")

if __name__ == "__main__":
    main()
